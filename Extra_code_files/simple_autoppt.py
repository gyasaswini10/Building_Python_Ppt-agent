#!/usr/bin/env python3
"""
Simple Working AutoPPT Agent

Minimal working version that creates PowerPoint presentations.
"""

import asyncio
import json
import os
import re
import sys
import time
from pathlib import Path
from typing import List, Dict, Any, Optional

try:
    from dotenv import load_dotenv
    import httpx
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError as e:
    print(f"Missing required package: {e}")
    print("Install with: pip install python-pptx httpx python-dotenv")
    sys.exit(1)

load_dotenv()

class SimpleAutoPPT:
    """Simple working AutoPPT agent."""
    
    def __init__(self, request: str, output_file: str = "presentation.pptx"):
        self.request = request
        self.output_file = output_file
        self.session_id = f"session_{int(time.time())}"
        
        print(f"🎯 Simple AutoPPT Agent")
        print(f"📝 Request: {request}")
        print(f"📁 Output: {output_file}")
    
    def extract_topic(self, request: str) -> str:
        """Extract main topic from request."""
        pattern1 = r'(?i)^(create|make|generate|build)\s+(a\s+)?(\d+-slide\s+)?presentation\s+(?:on|about|for)\s+(.+)$'
        match = re.match(pattern1, request.strip())
        if match:
            return match.group(3).strip()
        
        pattern2 = r'(?i)^presentation\s+(?:on|about|for)\s+(.+)$'
        match = re.match(pattern2, request.strip())
        if match:
            return match.group(1).strip()
        
        return request.strip()
    
    def create_simple_presentation(self, title: str) -> Dict[str, Any]:
        """Create a simple PowerPoint presentation."""
        try:
            self.presentation = Presentation()
            
            # Create title slide
            slide_layout = self.presentation.slide_layouts[0]
            title_slide = self.presentation.slides.add_slide(slide_layout)
            
            # Add title
            if title_slide.shapes.title:
                title_shape = title_slide.shapes.title
                title_shape.text = title
                title_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(36)
                title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 223, 128)
                title_shape.text_frame.paragraphs[0].runs[0].font.bold = True
            
            print(f"📊 Created presentation: {title}")
            return {"ok": True, "session_id": self.session_id, "title": title}
        except Exception as e:
            return {"ok": False, "error": str(e)}
    
    def add_content_slide(self, title: str, bullets: List[str]) -> Dict[str, Any]:
        """Add a content slide."""
        if not self.presentation:
            return {"ok": False, "error": "No presentation created"}
        
        try:
            slide_layout = self.presentation.slide_layouts[0]
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Add title
            title_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
            title_shape.text = title
            title_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(28)
            title_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 223, 128)
            title_shape.text_frame.paragraphs[0].runs[0].font.bold = True
            
            # Add bullet points
            for i, bullet in enumerate(bullets[:6]):
                bullet_shape = slide.shapes.add_textbox(Inches(1), Inches(2.5 + i * 0.8), Inches(8), Inches(0.6))
                bullet_shape.text = f"• {bullet}"
                bullet_shape.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
                bullet_shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
            
            print(f"📝 Added slide: {title}")
            return {"ok": True, "slide_index": len(self.presentation.slides) - 1}
        except Exception as e:
            return {"ok": False, "error": str(e)}
    
    def save_presentation(self, output_path: str) -> Dict[str, Any]:
        """Save the presentation."""
        if not self.presentation:
            return {"ok": False, "error": "No presentation to save"}
        
        try:
            output_dir = Path(output_path).parent
            output_dir.mkdir(parents=True, exist_ok=True)
            
            # Save without complex formatting
            self.presentation.save(output_path)
            
            file_size = Path(output_path).stat().st_size
            print(f"💾 Saved presentation: {output_path} ({file_size:,} bytes)")
            return {"ok": True, "output_path": output_path, "file_size": file_size}
        except Exception as e:
            return {"ok": False, "error": str(e)}
    
    async def execute(self) -> Dict[str, Any]:
        """Execute the presentation creation."""
        try:
            print("\n🚀 Starting Simple AutoPPT Agent")
            print("=" * 50)
            
            # Extract topic
            topic = self.extract_topic(self.request)
            print(f"📋 Topic extracted: {topic}")
            
            # Create presentation
            create_result = self.create_simple_presentation(f"Presentation on {topic}")
            if not create_result["ok"]:
                return {"ok": False, "error": create_result["error"]}
            
            # Generate content slides
            slide_content = [
                ("Introduction", [f"Overview of {topic}", f"Key concepts in {topic}", f"Importance of {topic}"]),
                ("Main Features", [f"Feature 1 of {topic}", f"Feature 2 of {topic}", f"Feature 3 of {topic}"]),
                ("Applications", [f"Application 1", f"Application 2", f"Application 3"]),
                ("Benefits", [f"Benefit 1", f"Benefit 2", f"Benefit 3"]),
                ("Conclusion", [f"Summary of {topic}", f"Future outlook", f"Key takeaways"])
            ]
            
            for title, bullets in slide_content:
                slide_result = self.add_content_slide(title, bullets)
                if not slide_result["ok"]:
                    print(f"❌ Failed to add slide {title}: {slide_result['error']}")
            
            # Save presentation
            save_result = self.save_presentation(self.output_file)
            
            if save_result["ok"]:
                print("\n🎉 Presentation created successfully!")
                print(f"📁 Location: {Path(self.output_file).absolute()}")
                print(f"📊 File size: {save_result['file_size']:,} bytes")
                return {
                    "ok": True,
                    "output_path": save_result["output_path"],
                    "file_size": save_result["file_size"],
                    "slides_created": len(self.presentation.slides)
                }
            else:
                return {"ok": False, "error": save_result["error"]}
                
        except Exception as e:
            print(f"❌ Workflow error: {e}")
            return {"ok": False, "error": str(e)}

def main():
    """Main entry point."""
    if len(sys.argv) < 2:
        print("🎯 Simple AutoPPT Agent")
        print("=" * 30)
        print("Usage:")
        print("  python simple_autoppt.py \"Create a 5-slide presentation on solar energy\"")
        print("  python simple_autoppt.py \"Make a presentation about AI\" --output \"ai.pptx\"")
        return
    
    request = sys.argv[1]
    output_file = "presentation.pptx"
    
    # Parse optional arguments
    for i, arg in enumerate(sys.argv[2:], 2):
        if arg == "--output" and i + 1 < len(sys.argv):
            output_file = sys.argv[i + 1]
        elif arg.startswith("--output="):
            output_file = arg.split("=", 1)[1]
    
    # Create and run agent
    agent = SimpleAutoPPT(request, output_file)
    
    async def run_agent():
        result = await agent.execute()
        if result["ok"]:
            print(f"\n✅ Success! Created {result['slides_created']} slides")
        else:
            print(f"\n❌ Failed: {result['error']}")
    
    print("\n🚀 Starting Simple AutoPPT Agent...")
    asyncio.run(run_agent())

if __name__ == "__main__":
    main()
