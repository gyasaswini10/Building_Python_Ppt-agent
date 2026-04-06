"""
Flask Web Server for MCP-PPT-Agent Frontend Integration
Single command solution that uses existing modular code files
"""

import os
import threading
import sys
import json
import asyncio
import uuid
import httpx
from pathlib import Path
from flask import Flask, request, jsonify, send_file
from flask_cors import CORS

# Add the current directory's Modular code folder to sys.path
sys.path.append(str(Path(__file__).parent / "Modular code"))

# Load environment variables from .env file
try:
    from dotenv import load_dotenv
    load_dotenv()
    print("✅ Environment variables loaded")
except ImportError:
    print("⚠️ python-dotenv not available, using defaults")

app = Flask(__name__)
CORS(app)

# Configuration
# Explicitly set the saving path as requested by the user
app.config['UPLOAD_FOLDER'] = r'C:\Users\gyasu\Desktop\CAlibo noww\ASSIGNMENT\Client\savingfolder_output'
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 16MB max file size

# Ensure savingfolder_output directory exists
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Global variables for session management
active_sessions = {}

# Try to import python-pptx for PowerPoint generation
try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
    from pptx.enum.text import MSO_AUTO_SIZE
    PPTX_AVAILABLE = True
    print("✅ python-pptx available - Full PowerPoint generation enabled")
except ImportError:
    PPTX_AVAILABLE = False
    print("⚠️ python-pptx not available - Running in demo mode")

# Initializing global state for MCP communication (Session IDs are ephemeral)
active_sessions = {}
# Thread-safe lock for session modifications
session_lock = threading.Lock()
# Correct modern import for the AI Brain and Modular Servers
try:
    import agent_ppt
    import ppt_mcp_server as ppt_module
    import research_mcp_server as research_module
    # Force reload to ensure constructor changes are picked up without restarting thread
    import importlib
    importlib.reload(agent_ppt)
    MODULAR_CODE_AVAILABLE = True
    print("✅ Localized AI Brain and Modular Servers RELOADED successfully")
except ImportError as e:
    print(f"⚠️ Modular code components not found in path: {e}")
    MODULAR_CODE_AVAILABLE = False

async def research_topic_robust(query: str, slide_title: str = "") -> dict:
    """Research function with localized AI brain for thread safety"""
    # --- STAGE 1: DYNAMIC AI RESEARCH (OpenRouter/HF) ---
    if MODULAR_CODE_AVAILABLE:
        try:
            # Force a localized reload to avoid any stale class definitions in memory
            import importlib
            importlib.reload(agent_ppt)
            
            # Pass two arguments (topic, path) to satisfy both old and new constructor versions
            brain = agent_ppt.AutonomousPresenter(query, "presentation.pptx")
            brain.setup_llm()
            
            print(f"🧪 AI Pipeline: Seeking Data for '{slide_title}' (Keys: {len(brain.openrouter_keys)})")
            
            search_query = f"{query} - {slide_title}" if slide_title else query
            res = await brain.ask_llm(search_query)
            
            if res.get("bullets") and len(res["bullets"]) >= 1:
                return {"ok": True, "points": res["bullets"][:6], "source": "LLM-Orchestrator"}
            else:
                print(f"⚠️ AI Pipeline: No points returned for '{slide_title}'. Trying Wikipedia fallback...")
        except Exception as e:
            print(f"⚠️ AI Pipeline Exception: {e}")

    # --- STAGE 2: MODULAR Wikipedia FETCH (Standard MCP) ---
    if MODULAR_CODE_AVAILABLE and hasattr(research_module, 'search_topic'):
        try:
            res = await research_module.search_topic(query, slide_title)
            if res and res.get("points"): return res
        except Exception: pass

    # --- STAGE 3: DIRECT Wikipedia FALLBACK ---
    try:
        q = query.strip()
        for suffix in (" facts biology", " facts", " overview"):
            if q.lower().endswith(suffix): q = q[: -len(suffix)].strip()
        wiki_title = q.replace(" ", "_")
        url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{wiki_title}"
        
        async with httpx.AsyncClient(timeout=10.0, follow_redirects=True) as client:
            resp = await client.get(url)
            if resp.status_code == 200:
                data = resp.json()
                extract = (data.get("extract", "") or "").strip()
                if extract:
                    sentences = [s.strip() + '.' for s in extract.split('. ') if len(s.strip()) > 20]
                    return {"ok": True, "points": sentences[:5], "source": "wikipedia-fallback"}
    except Exception: pass
    
    return {"ok": False, "error": "AI and Wiki research both failed", "points": []}

# PowerPoint creation functions
def create_presentation_simple(title: str = "Presentation") -> dict:
    """Create a new PowerPoint presentation"""
    if not PPTX_AVAILABLE:
        # Return demo session
        session_id = str(uuid.uuid4())
        active_sessions[session_id] = {
            "title": title,
            "slides": [],
            "demo": True
        }
        return {"ok": True, "session_id": session_id}
    
    try:
        pres = Presentation()
        slide = pres.slides.add_slide(pres.slide_layouts[0])
        
        # Theme: Midnight Navy background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 25, 45)
        
        # Title slide styling
        title_shape = slide.shapes.title
        title_shape.text = title
        for tr in title_shape.text_frame.paragraphs[0].runs:
            tr.font.color.rgb = RGBColor(76, 175, 80) # Vibrant Green
            tr.font.size = Pt(44)
            tr.font.bold = True
        
        subtitle = slide.placeholders[1]
        subtitle.text = "Generated by MCP-PPT-Agent"
        sr = subtitle.text_frame.paragraphs[0].runs[0]
        sr.font.color.rgb = RGBColor(76, 175, 80) # Vibrant Green
        sr.font.size = Pt(18)
        
        session_id = str(uuid.uuid4())
        active_sessions[session_id] = {
            "presentation": pres,
            "title": title,
            "demo": False
        }
        
        return {"ok": True, "session_id": session_id}
    except Exception as e:
        print(f"PPT creation error: {e}")
        return {"ok": False, "error": str(e)}

def add_slide_simple(session_id: str, slide_title: str, bullets: list) -> dict:
    """Add a slide to the presentation"""
    if session_id not in active_sessions:
        return {"ok": False, "error": "Invalid session_id"}
    
    session = active_sessions[session_id]
    
    if session.get("demo"):
        # Demo mode - just store slide data
        session["slides"].append({
            "title": slide_title,
            "bullets": bullets
        })
        return {"ok": True, "slides_total": len(session["slides"])}
    
    try:
        pres = session["presentation"]
        slide = pres.slides.add_slide(pres.slide_layouts[1])
        
        # Theme: Midnight Navy background
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(15, 25, 45)
        
        # Slide title
        slide.shapes.title.text = slide_title
        title_para = slide.shapes.title.text_frame.paragraphs[0]
        for tr in title_para.runs:
            tr.font.color.rgb = RGBColor(76, 175, 80) # Vibrant Green
            tr.font.bold = True
            tr.font.size = Pt(36)
        
        # Slide content
        body = slide.shapes.placeholders[1]
        tf = body.text_frame
        tf.clear()
        
        # Add bullets
        for i, bullet in enumerate(bullets[:5]):  # Limit to 5 bullets
            if i > 0:
                tf.add_paragraph()
            p = tf.paragraphs[i] if i < len(tf.paragraphs) else tf.add_paragraph()
            p.text = str(bullet).strip()
            p.level = 0
            for run in p.runs:
                run.font.size = Pt(18)
                run.font.color.rgb = RGBColor(76, 175, 80) # Vibrant Green
        
        return {"ok": True, "slides_total": len(pres.slides)}
    except Exception as e:
        print(f"Add slide error: {e}")
        return {"ok": False, "error": str(e)}

def save_presentation_simple(session_id: str, output_path: str) -> dict:
    """Save the presentation to disk"""
    if session_id not in active_sessions:
        return {"ok": False, "error": "Invalid session_id"}
    
    session = active_sessions[session_id]
    
    if session.get("demo"):
        # Demo mode - create a simple text file
        output_file = os.path.join(app.config['UPLOAD_FOLDER'], f"{session_id}_demo.txt")
        with open(output_file, 'w', encoding='utf-8') as f:
            f.write(f"Presentation: {session['title']}\n\n")
            for i, slide in enumerate(session["slides"], 1):
                f.write(f"Slide {i}: {slide['title']}\n")
                for bullet in slide["bullets"]:
                    f.write(f"• {bullet}\n")
                f.write("\n")
        return {"ok": True, "output_path": output_file}
    
    try:
        pres = session["presentation"]
        output_file = os.path.join(app.config['UPLOAD_FOLDER'], output_path)
        pres.save(output_file)
        return {"ok": True, "output_path": output_file}
    except Exception as e:
        print(f"Save error: {e}")
        return {"ok": False, "error": str(e)}

# Flask Routes
@app.route('/')
def index():
    """Serve the main frontend page"""
    return send_file('templates/index.html')

@app.route('/styles.css')
def styles():
    """Serve the CSS file"""
    return send_file('styles.css', mimetype='text/css')

@app.route('/script.js')
def script():
    """Serve the JavaScript file"""
    return send_file('script.js', mimetype='application/javascript')

@app.route('/api/health')
def health_check():
    """Health check endpoint"""
    return jsonify({
        "status": "healthy",
        "pptx_available": PPTX_AVAILABLE,
        "active_sessions": len(active_sessions)
    })

@app.route('/api/research-topic', methods=['POST'])
def research_topic():
    """Research a topic"""
    try:
        data = request.get_json()
        query = data.get('query', '')
        slide_title = data.get('slide_title', '')
        print(f"🔍 API: Researching '{query}' for slide '{slide_title}'")
        
        # Run async research
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        result = loop.run_until_complete(research_topic_robust(query, slide_title))
        loop.close()
        
        if result.get("ok"):
            print(f"✅ API: Research SUCCESS (Source: {result.get('source', 'unknown')})")
        else:
            print(f"⚠️ API: Research FAILED for '{query}'")
            
        return jsonify(result)
    except Exception as e:
        print(f"❌ API Error in research_topic: {e}")
        return jsonify({"ok": False, "error": str(e)})

@app.route('/api/create-presentation', methods=['POST'])
def create_presentation():
    """Create a new presentation"""
    try:
        data = request.get_json()
        title = data.get('title', 'Presentation')
        print(f"🎬 API: Creating presentation with title: {title}")
        
        # Try modular logic first
        if MODULAR_CODE_AVAILABLE and hasattr(ppt_module, 'create_pptx'):
            print(f"🚀 Using MODULAR ppt_module to create: {title}")
            result = ppt_module.create_pptx(title)
            if result.get("ok"):
                print(f"✅ Modular PPT created. Session: {result.get('session_id')}")
                return jsonify(result)
            else:
                print(f"⚠️ Modular PPT failed: {result.get('error')}. Falling back to internal...")
        
        # Fallback/Direct implementation
        print(f"🔄 Using INTERNAL simple_creation for: {title}")
        result = create_presentation_simple(title)
        if result.get("ok"):
            print(f"✅ Internal PPT created. Session: {result.get('session_id')}")
        else:
            print(f"❌ Internal PPT failed: {result.get('error')}")
        return jsonify(result)
    except Exception as e:
        print(f"❌ API Error in create_presentation: {e}")
        return jsonify({"ok": False, "error": str(e)})

@app.route('/api/add-slide', methods=['POST'])
def add_slide():
    """Add a slide to presentation"""
    try:
        data = request.get_json()
        session_id = data.get('session_id')
        slide_title = data.get('slide_title', '')
        bullets = data.get('bullets', [])
        
        if not session_id or not slide_title:
            return jsonify({"ok": False, "error": "session_id and slide_title are required"})
        
        # Try modular logic first
        if MODULAR_CODE_AVAILABLE and hasattr(ppt_module, 'add_slide'):
            print(f"🚀 Attempting MODULAR add_slide for session {session_id}...")
            result = ppt_module.add_slide(session_id, slide_title, bullets)
            if result.get("ok"): 
                print(f"✅ Modular add_slide SUCCESS for {slide_title}")
                return jsonify(result)
            
            # If it failed due to session mismatch, we MUST fall back
            print(f"⚠️ Modular add_slide MISSED or FAILED: {result.get('error')}. Checking internal sessions...")
        
        # Fallback to internal logic
        print(f"🔄 Attempting INTERNAL add_slide for session {session_id}...")
        result = add_slide_simple(session_id, slide_title, bullets)
        if result.get("ok"):
            print(f"✅ Internal add_slide SUCCESS for {slide_title}")
        else:
            print(f"❌ Internal add_slide FAILED for session {session_id}")
            
        return jsonify(result)
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)})

@app.route('/api/delete-slide', methods=['POST'])
def delete_slide():
    """Delete a slide from presentation"""
    try:
        data = request.get_json()
        session_id = data.get('session_id')
        slide_index = data.get('slide_index')
        
        if not session_id or slide_index is None:
            return jsonify({"ok": False, "error": "session_id and slide_index are required"})
        
        # Use modular ppt tool if available
        if MODULAR_CODE_AVAILABLE and hasattr(ppt_module, 'delete_slide'):
            result = ppt_module.delete_slide(session_id, slide_index)
            return jsonify(result)
            
        # Fallback to internal/simplified delete
        if session_id not in active_sessions:
            return jsonify({"ok": False, "error": "Invalid session_id"})
        
        session = active_sessions[session_id]
        
        if session.get("demo"):
            # Demo mode - remove from slide list
            if 0 <= slide_index < len(session["slides"]):
                session["slides"].pop(slide_index)
                return jsonify({"ok": True, "slides_total": len(session["slides"])})
            else:
                return jsonify({"ok": False, "error": "Invalid slide_index"})
        
        # Real PowerPoint mode
        try:
            pres = session["presentation"]
            if 0 <= slide_index < len(pres.slides):
                # Remove the slide (using python-pptx internal list manipulation)
                xml_slides = pres.slides._sldIdLst
                slides = list(xml_slides)
                xml_slides.remove(slides[slide_index])
                return jsonify({"ok": True, "slides_total": len(pres.slides)})
            else:
                return jsonify({"ok": False, "error": "Invalid slide_index"})
        except Exception as e:
            return jsonify({"ok": False, "error": str(e)})
            
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)})

@app.route('/api/get-ppt-info', methods=['POST'])
def get_ppt_info():
    """Get information about a presentation"""
    try:
        data = request.get_json()
        session_id = data.get('session_id')
        
        if not session_id:
            return jsonify({"ok": False, "error": "session_id is required"})
        
        # Use modular ppt tool if available
        if MODULAR_CODE_AVAILABLE and hasattr(ppt_module, 'get_ppt_info'):
            result = ppt_module.get_ppt_info(session_id)
            if result["ok"]:
                # Map modular structure to frontend expected structure
                # Frontend expects "slide_index" and "bullet_count"
                slides = []
                for s in result["slides"]:
                    slides.append({
                        "slide_index": s["index"],
                        "title": s["title"],
                        "bullet_count": 5 # Modular code doesn't store count, but we show the max
                    })
                return jsonify({
                    "ok": True,
                    "slides": slides,
                    "slides_total": result["slide_count"],
                    "title": "Presentation"
                })
            return jsonify(result)

        if session_id not in active_sessions:
            return jsonify({"ok": False, "error": "Invalid session_id"})
        
        session = active_sessions[session_id]
        
        if session.get("demo"):
            # Demo mode
            slides_info = []
            for i, slide in enumerate(session["slides"]):
                slides_info.append({
                    "slide_index": i,
                    "title": slide["title"],
                    "bullet_count": len(slide["bullets"])
                })
            
            return jsonify({
                "ok": True,
                "slides_total": len(session["slides"]),
                "slides": slides_info,
                "title": session.get("title", "Presentation")
            })
        
        # Real PowerPoint mode
        try:
            pres = session["presentation"]
            slides_info = []
            
            for i, slide in enumerate(pres.slides):
                slide_title = slide.shapes.title.text if slide.shapes.title else f"Slide {i+1}"
                bullet_count = 0
                
                if slide.shapes.placeholders:
                    try:
                        body = slide.shapes.placeholders[1]
                        bullet_count = len([p for p in body.text_frame.paragraphs if p.text.strip()])
                    except:
                        pass
                
                slides_info.append({
                    "slide_index": i,
                    "title": slide_title,
                    "bullet_count": bullet_count
                })
            
            return jsonify({
                "ok": True,
                "slides_total": len(pres.slides),
                "slides": slides_info,
                "title": session.get("title", "Presentation")
            })
            
        except Exception as e:
            return jsonify({"ok": False, "error": str(e)})
            
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)})

@app.route('/api/save-presentation', methods=['POST'])
def save_presentation():
    """Save presentation"""
    try:
        data = request.get_json()
        session_id = data.get('session_id')
        output_path = data.get('output_path', 'presentation.pptx')
        
        if not session_id:
            return jsonify({"ok": False, "error": "session_id is required"})
            
        print(f"💾 API: Saving session {session_id} to {output_path}")
        
        # NEW: Automatic Versioning Logic
        # Check if folder exists, then check for filename conflicts
        if not os.path.exists(app.config['UPLOAD_FOLDER']):
            os.makedirs(app.config['UPLOAD_FOLDER'])
            
        full_output_path = os.path.join(app.config['UPLOAD_FOLDER'], output_path)
        
        if os.path.exists(full_output_path):
            print(f"⚠️ Conflict! {output_path} exists. Indexing version...")
            base, ext = os.path.splitext(output_path)
            counter = 1
            while os.path.exists(full_output_path):
                counter += 1
                new_name = f"{base}_v{counter}{ext}"
                full_output_path = os.path.join(app.config['UPLOAD_FOLDER'], new_name)
            output_path = os.path.basename(full_output_path)
            print(f"✅ Versioning: New destination set to {output_path}")

        # Try modular save first if available
        if MODULAR_CODE_AVAILABLE and hasattr(ppt_module, 'save_presentation'):
            print(f"🚀 Using MODULAR save logic at {full_output_path}...")
            result = ppt_module.save_presentation(session_id, full_output_path)
        else:
            print(f"🔄 Using INTERNAL save logic at {full_output_path}...")
            # Internal save might ignore our full_path, so we use simpler call if needed
            result = save_presentation_simple(session_id, output_path)
            
        if result.get("ok"):
            print(f"✅ Presentation saved to: {result.get('output_path')}")
        else:
            print(f"❌ Save FAILED: {result.get('error')}")
            
        return jsonify(result)
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)})

@app.route('/api/download/<session_id>')
def download_presentation(session_id):
    """Download presentation file"""
    try:
        print(f"⬇️ API: Download request received for session: {session_id}")
        
        # Look for saved file matching session_id snippet
        files = os.listdir(app.config['UPLOAD_FOLDER'])
        print(f"📂 Scanning storage for matching file in: {app.config['UPLOAD_FOLDER']}")
        
        for filename in files:
            # We save with topic_sessionPrefix so we check if session_id or part of it is in name
            if session_id in filename or session_id[:8] in filename:
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                print(f"✅ Match Found! Serving: {filename}")
                return send_file(file_path, as_attachment=True, download_name=filename)
        
        print(f"❌ Download FAILED: No file matching {session_id[:8]} found in storage.")
        return jsonify({"error": "File not found"}), 404
    except Exception as e:
        print(f"❌ Download Route Error: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/generate-presentation', methods=['POST'])
def generate_presentation():
    """Complete presentation generation in one call"""
    try:
        data = request.get_json()
        topic = data.get('topic', '')
        title = data.get('title', f'Presentation on {topic}')
        
        if not topic:
            return jsonify({"ok": False, "error": "Topic is required"})
        
        # Create presentation - Try modular first
        if MODULAR_CODE_AVAILABLE and hasattr(ppt_module, 'create_pptx'):
            print(f"🚀 Automator: Creating MODULAR presentation: {title}")
            create_result = ppt_module.create_pptx(title)
        else:
            print(f"🔄 Automator: Creating INTERNAL simple presentation: {title}")
            create_result = create_presentation_simple(title)
            
        if not create_result.get("ok"):
            print(f"❌ Automator: Creation FAILED - {create_result.get('error')}")
            return jsonify(create_result)
        
        session_id = create_result["session_id"]
        print(f"✅ Automator: Session started: {session_id}")
        
        # Define slide structure (6-slide scientific hierarchy)
        slide_topics = [
            "Origins and Taxonomy",
            "Physiological & Structural Features", 
            "Biological Growth & Lifecycle",
            "Ecological & Environmental Roles",
            "Industrial & Societal Impact",
            "Future Research & Conservation"
        ]
        
        # Generate slides
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        
        # Initialize agent for LLM-backed research
        agent = None
        if MODULAR_CODE_AVAILABLE:
            try:
                agent = agent_module.AutonomousPresenter(topic, os.path.join(app.config['UPLOAD_FOLDER'], filename))
                agent.setup_llm()
            except Exception as e:
                print(f"Agent initialization error: {e}")

        for slide_topic in slide_topics:
            bullets = []
            # Try to get bullets from Agent LLM first
            if agent:
                try:
                    # Pass only the specific slide heading to the LLM agent
                    # The prompt context is now handled inside ask_llm for better factual results
                    result = loop.run_until_complete(agent.ask_llm(f"{topic} - {slide_topic}"))
                    bullets = result.get("bullets", [])
                except Exception as e:
                    print(f"Agent LLM error: {e}")
            
            # Fallback to robust research (Wikipedia) if LLM fails or returns nothing
            if not bullets:
                research_result = loop.run_until_complete(research_topic_robust(topic, slide_topic))
                if research_result.get("ok") and research_result.get("points"):
                    bullets = research_result["points"][:5]
            
            # Final fallback if both fail
            if not bullets:
                bullets = [f"Detailed analysis of {slide_topic}", 
                          "Comprehensive research findings",
                          "Key takeaway and structural overview",
                          "Practical implications and results",
                          "Future considerations for this area"]
            
            # Add slide - Try modular first, then internal
            success = False
            print(f"📽️ Processing slide: '{slide_topic}' (Bullets: {len(bullets)})")
            
            if MODULAR_CODE_AVAILABLE and hasattr(ppt_module, 'add_slide'):
                print(f"🚀 Automator: Attempting MODULAR add_slide for {slide_topic}")
                res = ppt_module.add_slide(session_id, slide_topic, bullets)
                if res.get("ok"): 
                    print(f"✅ Automator: Modular add SUCCESS for index {res.get('slides_total', 0)-1}")
                    success = True
            
            if not success:
                print(f"🔄 Automator: Falling back to INTERNAL add for {slide_topic}")
                internal_res = add_slide_simple(session_id, slide_topic, bullets)
                print(f"✅ Automator: Internal add COMPLETE. Total slides: {internal_res.get('slides_total', 'unknown')}")
        
        loop.close()
        
        # Save presentation - include session_id in filename for download route
        topic_clean = topic.replace(' ', '_').replace('/', '_')
        filename = f"{topic_clean}_{session_id[:8]}.pptx"
        full_filename = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        
        if MODULAR_CODE_AVAILABLE and hasattr(ppt_module, 'save_presentation'):
            save_result = ppt_module.save_presentation(session_id, full_filename)
        else:
            save_result = save_presentation_simple(session_id, filename)
        
        if save_result["ok"]:
            return jsonify({
                "ok": True,
                "session_id": session_id,
                "filename": filename,
                "slides_created": len(slide_topics)
            })
        else:
            return jsonify(save_result)
            
    except Exception as e:
        return jsonify({"ok": False, "error": str(e)})

if __name__ == '__main__':
    print("🚀 Starting MCP-PPT-Agent Web Server")
    print(f"📁 Output folder: {app.config['UPLOAD_FOLDER']}")
    print(f"🌐 Server will be available at: http://localhost:5000")
    print(f"📊 PowerPoint Generation: {'✅ Available' if PPTX_AVAILABLE else '⚠️ Demo Mode'}")
    
    app.run(host='0.0.0.0', port=5000, debug=True)
