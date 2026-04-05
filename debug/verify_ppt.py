from pptx import Presentation
import sys
import os

def read_pptx(path):
    if not os.path.exists(path):
        print(f"File not found: {path}")
        return
    prs = Presentation(path)
    print(f"Reading {path}: {len(prs.slides)} slides found.")
    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i+1} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    print(f"[{shape.name}]: {paragraph.text}")

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python verify_ppt.py <path>")
    else:
        read_pptx(sys.argv[1])
