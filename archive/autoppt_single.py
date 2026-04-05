#!/usr/bin/env python3
"""
AutoPPT Single-File Agent

Complete PowerPoint creation system that runs with a single command.
No separate MCP servers needed - everything is integrated.

Created by: Yasaswini
Course: AI Agents & MCP Architecture
"""

import asyncio
import json
import os
import re
import sys
import time
from pathlib import Path
from typing import List, Dict, Any, Optional

# Add environment setup
sys.path.append(str(Path(__file__).parent))

try:
    from dotenv import load_dotenv
    import httpx
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_AUTO_SIZE
except ImportError as e:
    print(f"Missing required package: {e}")
    print("Install with: pip install -r requirements.txt")
    sys.exit(1)

# Load environment variables
load_dotenv()

class AutoPPTSingle:
    """Single-file AutoPPT agent with all functionality integrated."""
    
    def __init__(self, request: str, output_file: str = "presentation.pptx"):
        """Initialize the AutoPPT agent."""
        self.request = request
        self.output_file = output_file
        self.presentation = None
        self.session_id = f"session_{int(time.time())}"
        
        # API keys management
        self.openrouter_keys = self._get_api_keys("OPENROUTER_API_KEY")
        self.hf_tokens = self._get_api_keys("HF_TOKENS")
        self.current_openrouter_key = 0
        self.current_hf_token = 0
        
        print(f"🎯 AutoPPT Agent Initialized")
        print(f"📝 Request: {request}")
        print(f"📁 Output: {output_file}")
        print(f"🔑 API Keys: {len(self.openrouter_keys)} OpenRouter, {len(self.hf_tokens)} HuggingFace")
    
    def _get_api_keys(self, prefix: str) -> List[str]:
        """Extract API keys from environment variables."""
        keys = []
        
        # Try numbered keys first
        for i in range(1, 10):
            key = os.getenv(f"{prefix}_{i}")
            if key:
                keys.append(key)
        
        # Try comma-separated
        if not keys:
            combined = os.getenv(prefix)
            if combined:
                keys = [k.strip() for k in combined.split(",")]
        
        return [k for k in keys if k.strip()]
    
    def _get_available_openrouter_key(self) -> Optional[str]:
        """Get next available OpenRouter key."""
        for i in range(len(self.openrouter_keys)):
            key_idx = (self.current_openrouter_key + i) % len(self.openrouter_keys)
            key = self.openrouter_keys[key_idx]
            if key not in getattr(self, '_failed_openrouter_keys', set()):
                self.current_openrouter_key = key_idx
                return key
        return None
    
    def _get_available_hf_token(self) -> Optional[str]:
        """Get next available HuggingFace token."""
        for i in range(len(self.hf_tokens)):
            token_idx = (self.current_hf_token + i) % len(self.hf_tokens)
            token = self.hf_tokens[token_idx]
            if token not in getattr(self, '_failed_hf_tokens', set()):
                self.current_hf_token = token_idx
                return token
        return None
    
    async def _try_openrouter(self, prompt: str) -> Optional[str]:
        """Try to get response from OpenRouter API."""
        key = self._get_available_openrouter_key()
        if not key:
            return None
        
        try:
            async with httpx.AsyncClient(timeout=30.0) as client:
                response = await client.post(
                    "https://openrouter.ai/api/v1/chat/completions",
                    headers={
                        "Authorization": f"Bearer {key}",
                        "Content-Type": "application/json"
                    },
                    json={
                        "model": "anthropic/claude-3.5-sonnet",
                        "messages": [{"role": "user", "content": prompt}],
                        "max_tokens": 1000
                    }
                )
                
                if response.status_code == 200:
                    data = response.json()
                    return data["choices"][0]["message"]["content"]
                else:
                    print(f"OpenRouter error: {response.status_code}")
                    return None
                    
        except Exception as e:
            print(f"OpenRouter exception: {e}")
            return None
    
    async def _try_huggingface(self, prompt: str) -> Optional[str]:
        """Try to get response from HuggingFace API."""
        token = self._get_available_hf_token()
        if not token:
            return None
        
        try:
            async with httpx.AsyncClient(timeout=30.0) as client:
                response = await client.post(
                    "https://api-inference.huggingface.co/models/microsoft/DialoGPT-medium",
                    headers={
                        "Authorization": f"Bearer {token}",
                        "Content-Type": "application/json"
                    },
                    json={
                        "inputs": prompt,
                        "parameters": {
                            "max_length": 500,
                            "temperature": 0.7
                        }
                    }
                )
                
                if response.status_code == 200:
                    data = response.json()
                    if isinstance(data, list) and data:
                        return data[0].get("generated_text", "")
                    return str(data)
                else:
                    print(f"HuggingFace error: {response.status_code}")
                    return None
                    
        except Exception as e:
            print(f"HuggingFace exception: {e}")
            return None
    
    async def _generate_content(self, prompt: str) -> Optional[str]:
        """Generate content using available APIs."""
        # Try OpenRouter first
        result = await self._try_openrouter(prompt)
        if result:
            print("✅ OpenRouter API succeeded")
            return result
        
        # Fallback to HuggingFace
        result = await self._try_huggingface(prompt)
        if result:
            print("✅ HuggingFace API succeeded")
            return result
        
        print("❌ All APIs failed")
        return None
    
    def create_presentation(self, title: str) -> Dict[str, Any]:
        """Create a new PowerPoint presentation."""
        try:
            self.presentation = Presentation()
            
            # Set up slide size
            self.presentation.slide_width = Inches(10)
            self.presentation.slide_height = Inches(7.5)
            
            print(f"📊 Created presentation: {title}")
            return {
                "ok": True,
                "session_id": self.session_id,
                "title": title
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}
    
    def add_slide(self, slide_title: str, bullets: List[str]) -> Dict[str, Any]:
        """Add a slide with title and bullet points."""
        if not self.presentation:
            return {"ok": False, "error": "No presentation created"}
        
        try:
            from pptx.util import Inches
            
            # Create slide
            slide_layout = self.presentation.slide_layouts[0]
            slide = self.presentation.slides.add_slide(slide_layout)
            
            # Add title
            if slide.shapes.title:
                title_box = slide.shapes.title
                title_box.text = slide_title
                title_box.text_frame.font.size = Pt(28)
                title_box.text_frame.font.color.rgb = RGBColor(255, 223, 128)  # Gold color
                title_box.text_frame.font.bold = True
            
            # Add bullet points
            left = Inches(1.5)
            top = Inches(2.5)
            
            for i, bullet in enumerate(bullets[:6]):  # Limit to 6 bullets
                # Create bullet point
                txBox = slide.shapes.add_textbox(left, top + Inches(i * 0.8), Inches(8), Inches(0.6))
                tf = txBox.text_frame
                tf.text = f"• {bullet}"
                tf.font.size = Pt(18)
                tf.font.color.rgb = RGBColor(255, 255, 255)  # White color
            
            print(f"📝 Added slide: {slide_title}")
            return {
                "ok": True,
                "slide_index": len(self.presentation.slides) - 1,
                "title": slide_title
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}
    
    def save_presentation(self, output_path: str) -> Dict[str, Any]:
        """Save the presentation to file."""
        if not self.presentation:
            return {"ok": False, "error": "No presentation to save"}
        
        try:
            # Ensure output directory exists
            output_dir = Path(output_path).parent
            output_dir.mkdir(parents=True, exist_ok=True)
            
            # Apply dark theme to all slides
            for slide in self.presentation.slides:
                if hasattr(slide.background, 'fill'):
                    slide.background.fill.solid()
                    slide.background.fill.solid_color.rgb = RGBColor(15, 25, 45)  # Navy blue
            
            # Save the presentation
            self.presentation.save(output_path)
            
            file_size = Path(output_path).stat().st_size
            print(f"💾 Saved presentation: {output_path} ({file_size:,} bytes)")
            return {
                "ok": True,
                "output_path": output_path,
                "file_size": file_size
            }
        except Exception as e:
            return {"ok": False, "error": str(e)}
    
    async def research_topic(self, query: str) -> List[str]:
        """Research a topic using Wikipedia API."""
        try:
            # Clean up the query
            topic = re.sub(r'\s+(Stage|Part|Phase|Process|Overview|Introduction|Cycle)\b', '', query, flags=re.IGNORECASE).strip()
            topic = re.sub(r'\d+', '', topic).strip() or query
            
            # Search Wikipedia
            search_url = f"https://en.wikipedia.org/api/rest_v1/page/summary/{topic.replace(' ', '_')}"
            
            async with httpx.AsyncClient(timeout=10.0) as client:
                response = await client.get(search_url)
                
                if response.status_code == 200:
                    data = response.json()
                    extract = data.get("extract", "")
                    
                    if extract and len(extract) > 50:
                        # Split into sentences and create bullet points
                        sentences = re.split(r'[.!?]+', extract)
                        bullets = []
                        
                        for sentence in sentences[:6]:  # Limit to 6 bullets
                            sentence = sentence.strip()
                            if len(sentence) > 10:
                                bullets.append(sentence)
                        
                        print(f"🔍 Researched topic: {topic}")
                        return bullets
                    else:
                        # Fallback content
                        return [
                            f"Introduction to {topic}",
                            f"Key concepts in {topic}",
                            f"Applications of {topic}",
                            f"Future of {topic}",
                            f"Conclusion on {topic}"
                        ]
                else:
                    print(f"Wikipedia error: {response.status_code}")
                    return self._get_fallback_content(query)
                    
        except Exception as e:
            print(f"Research error: {e}")
            return self._get_fallback_content(query)
    
    def _get_fallback_content(self, topic: str) -> List[str]:
        """Generate fallback content when research fails."""
        return [
            f"Overview of {topic}",
            f"Main characteristics of {topic}",
            f"Important aspects of {topic}",
            f"Applications and uses",
            f"Future developments",
            f"Summary and conclusions"
        ]
    
    def extract_topic(self, request: str) -> str:
        """Extract main topic from user request."""
        # Pattern matching for "Create a N-slide presentation on TOPIC"
        import re
        
        pattern1 = r'(?i)^(create|make|generate|build)\s+(a\s+)?(\d+-slide\s+)?presentation\s+(?:on|about|for)\s+(.+)$'
        match = re.match(pattern1, request.strip())
        
        if match:
            return match.group(3).strip()
        
        # Pattern matching for "presentation on TOPIC"
        pattern2 = r'(?i)^presentation\s+(?:on|about|for)\s+(.+)$'
        match = re.match(pattern2, request.strip())
        
        if match:
            return match.group(1).strip()
        
        return request.strip()
    
    async def execute(self) -> Dict[str, Any]:
        """Execute the complete presentation creation workflow."""
        try:
            print("\n🚀 Starting AutoPPT Agent Workflow")
            print("=" * 50)
            
            # Step 1: Extract topic
            topic = self.extract_topic(self.request)
            print(f"📋 Topic extracted: {topic}")
            
            # Step 2: Research content
            print("\n🔍 Researching content...")
            research_results = await self.research_topic(topic)
            
            # Step 3: Create presentation
            print("\n📊 Creating presentation...")
            create_result = self.create_presentation(f"Presentation on {topic}")
            
            if not create_result["ok"]:
                return {"ok": False, "error": create_result["error"]}
            
            # Step 4: Generate slide content
            print("\n🤖 Generating slide content...")
            
            # Generate content for each slide
            slide_titles = [
                "Introduction",
                "Key Concepts", 
                "Applications",
                "Benefits",
                "Challenges",
                "Conclusion"
            ]
            
            for i, title in enumerate(slide_titles):
                if i >= len(research_results):
                    # Generate content using AI
                    prompt = f"Generate 3-4 bullet points about {title} for a presentation on {topic}. Keep each point concise and informative."
                    content = await self._generate_content(prompt)
                    
                    if content:
                        # Extract bullet points from AI response
                        bullets = re.split(r'[\n•*-]', content)
                        bullets = [b.strip() for b in bullets if b.strip() and len(b.strip()) > 5]
                    else:
                        bullets = research_results[i:i+1] if i < len(research_results) else self._get_fallback_content(title)
                else:
                    bullets = research_results[i:i+1]
                
                # Add slide
                slide_result = self.add_slide(title, bullets[:6])  # Limit to 6 bullets
                if not slide_result["ok"]:
                    print(f"❌ Failed to add slide {title}: {slide_result['error']}")
            
            # Step 5: Save presentation
            print("\n💾 Saving presentation...")
            save_result = self.save_presentation(self.output_file)
            
            if save_result["ok"]:
                print("\n🎉 Presentation created successfully!")
                print(f"📁 Location: {Path(self.output_file).absolute()}")
                print(f"📊 File size: {save_result['file_size']:,} bytes")
                return {
                    "ok": True,
                    "output_path": save_result["output_path"],
                    "file_size": save_result["file_size"],
                    "slides_created": len(self.presentation.slides)
                }
            else:
                return {"ok": False, "error": save_result["error"]}
                
        except Exception as e:
            print(f"❌ Workflow error: {e}")
            return {"ok": False, "error": str(e)}

def main():
    """Main entry point for single-file AutoPPT agent."""
    if len(sys.argv) < 2:
        print("🎯 AutoPPT Single-File Agent")
        print("=" * 40)
        print("Usage:")
        print("  python autoppt_single.py \"Create a 5-slide presentation on solar energy\"")
        print("  python autoppt_single.py \"Make a presentation about AI\" --output \"ai.pptx\"")
        print("\nOptions:")
        print("  --output FILE    Output filename (default: presentation.pptx)")
        print("  --slides NUM     Number of slides (default: 6)")
        print("  --debug          Enable debug output")
        return
    
    # Parse arguments
    request = sys.argv[1]
    output_file = "presentation.pptx"
    
    # Parse optional arguments
    for i, arg in enumerate(sys.argv[2:], 2):
        if arg == "--output" and i + 1 < len(sys.argv):
            output_file = sys.argv[i + 1]
        elif arg.startswith("--output="):
            output_file = arg.split("=", 1)[1]
        elif arg == "--debug":
            os.environ["DEBUG"] = "1"
            print("🐛 Debug mode enabled")
    
    # Create and run agent
    agent = AutoPPTSingle(request, output_file)
    
    async def run_agent():
        result = await agent.execute()
        if result["ok"]:
            print(f"\n✅ Success! Created {result['slides_created']} slides")
        else:
            print(f"\n❌ Failed: {result['error']}")
    
    print("\n🚀 Starting AutoPPT Agent...")
    asyncio.run(run_agent())

if __name__ == "__main__":
    main()
